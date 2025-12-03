# ============================================================
#  Full Analytics Backend - Gemini + PDF + PPTX + Chatbot
# ============================================================

from flask import Flask, request, jsonify, send_file
from flask_cors import CORS
import pandas as pd
import json
import os
from datetime import datetime
import google.generativeai as genai
from werkzeug.utils import secure_filename
import io
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from pptx import Presentation
from pptx.util import Inches, Pt
from dotenv import load_dotenv

# -------------------------------------------------------------
# INIT
# -------------------------------------------------------------
app = Flask(__name__)
CORS(app, resources={r"/*": {"origins": "*"}}, supports_credentials=True)

UPLOAD_FOLDER = "uploads"
ALLOWED_EXTENSIONS = {"csv"}

os.makedirs(UPLOAD_FOLDER, exist_ok=True)

load_dotenv()

# -------------------------------------------------------------
# GEMINI CONFIG
# -------------------------------------------------------------

API_KEY = os.getenv("GOOGLE_API_KEY") or os.getenv("GEMINI_API_KEY")

if not API_KEY:
    raise RuntimeError("❌ No API key found! Add GOOGLE_API_KEY=... to your .env")

print("DEBUG: API Loaded =", API_KEY[:10] + "******")

genai.configure(api_key=API_KEY)

# ⭐ WORKING MODELS ⭐
GEMINI_MODEL_ID = "gemini-2.5-pro"   # guaranteed compatible


# -------------------------------------------------------------
# HELPERS
# -------------------------------------------------------------

def allowed_file(filename):
    return "." in filename and filename.rsplit(".", 1)[1].lower() in ALLOWED_EXTENSIONS


def analyze_dataframe(df):
    return {
        "total_rows": len(df),
        "total_columns": len(df.columns),
        "column_names": list(df.columns),
        "dtypes": df.dtypes.astype(str).to_dict(),
        "missing_values": df.isnull().sum().to_dict(),
        "numeric_summary": df.describe().to_dict()
            if len(df.select_dtypes(include=["number"]).columns) > 0
            else {},
        "sample_rows": df.head(5).to_dict("records")
    }


# -------------------------------------------------------------
# FALLBACK REPORT (NO AI)
# -------------------------------------------------------------

def generate_basic_report(summary):
    kpis = [
        f"Total rows: {summary['total_rows']}",
        f"Total columns: {summary['total_columns']}"
    ]

    for col, stats in list(summary["numeric_summary"].items())[:4]:
        kpis.append(f"Average {col}: {round(stats['mean'], 2)}")

    return {
        "data_understanding": f"Dataset has {summary['total_rows']} rows and {summary['total_columns']} columns.",
        "kpis": kpis,
        "trend_analysis": "Trend analysis unavailable (fallback mode).",
        "anomaly_detection": "Anomaly analysis unavailable.",
        "correlation_analysis": "Correlation unavailable.",
        "insights": ["Insights unavailable (fallback mode)."],
        "executive_summary": "Fallback summary – AI unavailable.",
        "business_recommendations": ["Re-run when AI is available."],
        "slide_deck": {},
        "chatbot_format": {}
    }


# -------------------------------------------------------------
# AI REPORT GENERATOR
# -------------------------------------------------------------

def generate_report_with_gemini(summary, columns, sample):
    prompt = f"""
You are a senior analytics engine. Generate a FULL analytics report.

SUMMARY:
{json.dumps(summary, indent=2)}

COLUMNS:
{columns}

SAMPLES:
{json.dumps(sample[:3], indent=2)}

Return ONLY JSON in this format:
{{
 "data_understanding": "",
 "kpis": [],
 "trend_analysis": "",
 "anomaly_detection": "",
 "correlation_analysis": "",
 "insights": [],
 "executive_summary": "",
 "business_recommendations": [],
 "slide_deck": {{
    "slide_1_title": "",
    "slide_2_kpis": [],
    "slide_3_charts": "",
    "slide_4_insights": [],
    "slide_5_anomalies": [],
    "slide_6_recommendations": [],
    "slide_7_summary": ""
 }},
 "chatbot_format": {{
    "kpis": [],
    "insights": [],
    "summary": "",
    "recommendations": [],
    "anomalies": [],
    "trend_analysis": "",
    "columns": []
 }}
}}
"""

    try:
        model = genai.GenerativeModel(GEMINI_MODEL_ID)
        response = model.generate_content(prompt)

        text = response.text
        first = text.find("{")
        last = text.rfind("}")
        json_data = json.loads(text[first:last + 1])

        return json_data

    except Exception as e:
        print("AI ERROR:", e)
        return generate_basic_report(summary)


# -------------------------------------------------------------
# ROUTES
# -------------------------------------------------------------

@app.route("/api/health")
def health():
    return jsonify({"status": "running", "model": GEMINI_MODEL_ID})


@app.route("/api/upload", methods=["POST"])
def upload_file():
    if "file" not in request.files:
        return jsonify({"error": "No file"}), 400

    file = request.files["file"]
    if file.filename == "":
        return jsonify({"error": "Empty filename"}), 400

    if not allowed_file(file.filename):
        return jsonify({"error": "Only CSV allowed"}), 400

    filename = secure_filename(file.filename)
    path = os.path.join(UPLOAD_FOLDER, filename)
    file.save(path)

    df = pd.read_csv(path)
    summary = analyze_dataframe(df)

    report = generate_report_with_gemini(summary, summary["column_names"], summary["sample_rows"])

    report["metadata"] = {
        "filename": filename,
        "timestamp": datetime.now().isoformat(),
        "rows": summary["total_rows"],
        "columns": summary["total_columns"],
        "model_used": GEMINI_MODEL_ID
    }

    return jsonify(report)


# -------------------------------------------------------------
# CHAT ENDPOINT
# -------------------------------------------------------------

@app.route("/api/chat", methods=["POST"])
def chat():
    data = request.json
    question = data.get("question", "")
    context = data.get("context", {})

    prompt = f"""
Answer this only using this context:

{json.dumps(context, indent=2)}

Question: {question}
"""

    try:
        model = genai.GenerativeModel(GEMINI_MODEL_ID)
        res = model.generate_content(prompt)
        return jsonify({"answer": res.text})
    except:
        return jsonify({"answer": "AI unavailable"})


# -------------------------------------------------------------
# PDF EXPORT
# -------------------------------------------------------------

@app.route("/api/export/pdf", methods=["POST"])
def export_pdf():
    data = request.json

    buffer = io.BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=letter)
    styles = getSampleStyleSheet()
    story = []

    def add_section(title, content):
        story.append(Paragraph(f"<b>{title}</b>", styles["Heading2"]))
        if isinstance(content, list):
            for c in content:
                story.append(Paragraph(f"• {c}", styles["Normal"]))
        else:
            story.append(Paragraph(str(content), styles["Normal"]))
        story.append(Spacer(1, 0.2 * inch))

    story.append(Paragraph("<b>Weekly Analytics Report</b>", styles["Title"]))
    add_section("Executive Summary", data.get("executive_summary", ""))
    add_section("KPIs", data.get("kpis", []))
    add_section("Insights", data.get("insights", []))
    add_section("Recommendations", data.get("business_recommendations", []))

    doc.build(story)
    buffer.seek(0)

    return send_file(
        buffer,
        download_name="analytics_report.pdf",
        as_attachment=True,
        mimetype="application/pdf"
    )


# -------------------------------------------------------------
# PPTX EXPORT (CORS FIXED)
# -------------------------------------------------------------

@app.route("/api/export/pptx", methods=["POST", "OPTIONS"])
def export_pptx():
    if request.method == "OPTIONS":
        res = jsonify({"status": "ok"})
        res.headers.add("Access-Control-Allow-Origin", "*")
        res.headers.add("Access-Control-Allow-Headers", "*")
        res.headers.add("Access-Control-Allow-Methods", "POST, OPTIONS")
        return res

    data = request.json
    prs = Presentation()

    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Weekly Analytics Report"
    slide.placeholders[1].text = "Generated by Automated Insights Engine"

    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Executive Summary"
    slide.placeholders[1].text = data.get("executive_summary", "")

    # Slide 3
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "KPIs"
    slide.placeholders[1].text = "\n".join([f"• {k}" for k in data.get("kpis", [])])

    # Slide 4
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Key Insights"
    slide.placeholders[1].text = "\n".join([f"• {i}" for i in data.get("insights", [])])

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)

    response = send_file(
        buffer,
        download_name="analytics_report.pptx",
        as_attachment=True,
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
    response.headers.add("Access-Control-Allow-Origin", "*")
    return response


# -------------------------------------------------------------
# RUN SERVER
# -------------------------------------------------------------

if __name__ == "__main__":
    app.run(debug=True, port=5000)
